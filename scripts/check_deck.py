#!/usr/bin/env python3
"""スライド規約の機械チェック（PPTX / HTML 共通）。

使い方:
  python3 check_deck.py out.pptx
  python3 check_deck.py deck.html

正典: references/slide-rules.md
終了コード: FAIL があれば 1。
"""
import re
import sys
import zipfile
from pathlib import Path

# 自ブランドで禁止するレガシー色があればここに列挙（HEX 6桁・#なし）
OLD_COLORS = []
TITLE_MAX = 40
EMU_W, EMU_H = 12192000, 6858000

fails, warns = [], []


def fail(msg):
    fails.append(msg)


def warn(msg):
    warns.append(msg)


STRICT_LEN = True  # PPTX: 1行40字。HTMLは2行まで許容(>60でFAIL)

# 表記ゆれ代表ペア（slide-rules §7.6: 1資料1用語）。両方の表記が同一資料に現れたら WARN。
# (ラベルA, パターンA, ラベルB, パターンB)
TERM_VARIANTS = [
    ("メモリー", r"メモリー", "メモリ", r"メモリ(?!ー)(?!・CPU)"),  # 計算機の「メモリ・CPU」は除外
    ("紐付", r"紐付", "紐づ", r"紐づ"),
    ("ユーザー", r"ユーザー", "利用者", r"利用者"),
    ("アセット", r"アセット", "資産", r"資産"),
    ("フォルダー", r"フォルダー", "フォルダ", r"フォルダ(?!ー)"),
    ("サーバー", r"サーバー", "サーバ", r"サーバ(?!ー)"),
    ("メンバー", r"メンバー", "メンバ", r"メンバ(?!ー)"),
    ("コンピューター", r"コンピューター", "コンピュータ", r"コンピュータ(?!ー)"),
    ("問い合わせ", r"問い合わせ", "問合せ", r"問合せ"),
]


# AI臭ワード（slide-rules §7.9 / references/ai-smell-lexicon.md）。高確度語のみ WARN。
AI_SMELL_WORDS = [
    "まさに", "非常に", "極めて", "圧倒的", "画期的", "革新的", "次世代の",
    "過言ではありません", "に他なりません", "シームレス", "シナジー", "ソリューション",
    "エンドツーエンド", "ブラッシュアップ", "付加価値",
    "寄り添い", "伴走し", "二人三脚", "さらなる高みへ", "邁進",
    "昨今", "変化の激しい", "という点において", "の観点から",
    "させていただきます", "いただけますと幸いです",
    "と言えるでしょう", "と考えられます", "することが可能です",
]


def check_ai_smell(pages):
    """slide-rules §7.9: AI臭の高確度語を検出（WARN。文脈上正当なら目視で無視してよい）"""
    hits = {}
    for i, txt in pages:
        found = [w for w in AI_SMELL_WORDS if w in txt]
        if found:
            hits[i] = found
    if hits:
        detail = "、".join(f"p{i}:「{'/'.join(ws[:3])}」" for i, ws in sorted(hits.items())[:6])
        warn(f"AI臭ワード検出: {detail}（§7.9 / ai-smell-lexicon.md。素の動詞・直球の言い方に置き換え）")
    dash_pages = sorted({i for i, txt in pages if " — " in txt or "—" in txt})
    if dash_pages:
        warn(f"ダッシュ「 — 」連結 p{dash_pages}（AI文体の典型。句点・「：」・括弧に置き換え。§7.9）")


def check_terms(pages):
    """pages: [(idx, text), ...] 資料全体で両方の表記が出たら WARN（意味が別なら目視で無視してよい）"""
    for la, pa, lb, pb in TERM_VARIANTS:
        hits_a = sorted({i for i, t in pages if re.search(pa, t)})
        hits_b = sorted({i for i, t in pages if re.search(pb, t)})
        if hits_a and hits_b:
            warn(f"表記ゆれ疑い: 「{la}」p{hits_a} と「{lb}」p{hits_b} が混在（§7.6 1資料1用語。別概念なら可・目視確認）")

def check_title(idx, title):
    t = title.strip()
    if not t:
        warn(f"p{idx}: タイトルが空（表紙/扉なら可）")
        return
    if re.search(r"(です|ます|でした|ました)[。．.]?$", t):
        fail(f"p{idx}: タイトルがですます調 → 体言止めに: 「{t}」")
    if STRICT_LEN and len(t) > TITLE_MAX:
        fail(f"p{idx}: タイトル {len(t)} 字（>{TITLE_MAX}・PPTXは1行）: 「{t}」")
    elif len(t) > 60:
        fail(f"p{idx}: タイトル {len(t)} 字（>60）: 「{t}」")
    elif len(t) > TITLE_MAX:
        warn(f"p{idx}: タイトル {len(t)} 字（2行になる想定。泣き別れしない改行位置か目視確認）: 「{t}」")
    if re.match(r"^(Step|STEP|ステップ)\s*\d", t):
        fail(f"p{idx}: タイトルに Step 連結（タグチップで表現）: 「{t}」")
    if re.search(r"^(この|その|ここまで)", t):
        fail(f"p{idx}: 他スライド参照語で始まるタイトル: 「{t}」")
    if t.count("（") + t.count("(") >= 2:
        warn(f"p{idx}: タイトルに丸括弧が多い: 「{t}」")


# ---------------------------------------------------------------- PPTX
def check_pptx(path):
    try:
        from pptx import Presentation
        from pptx.util import Emu
    except ImportError:
        sys.exit("python-pptx が必要: pip3 install python-pptx")
    prs = Presentation(path)
    if abs(prs.slide_width - EMU_W) > 2000 or abs(prs.slide_height - EMU_H) > 2000:
        fail(f"スライドサイズ {prs.slide_width}x{prs.slide_height} ≠ 16:9 {EMU_W}x{EMU_H}")

    titles = []
    term_pages = []
    with zipfile.ZipFile(path) as z:
        names = [n for n in z.namelist() if n.startswith("ppt/slides/slide") and n.endswith(".xml")]
        for n in sorted(names, key=lambda s: int(re.search(r"slide(\d+)", s).group(1))):
            xml = z.read(n).decode("utf8", "ignore")
            idx = int(re.search(r"slide(\d+)", n).group(1))
            term_pages.append((idx, " ".join(re.findall(r"<a:t>(.*?)</a:t>", xml, re.S))))
            # 角丸（高さ 0.4in=365760 EMU 以上の図形のみ）
            rr = 0
            for m in re.finditer(r"<p:sp>.*?</p:sp>", xml, re.S):
                sp = m.group(0)
                if 'prst="roundRect"' in sp:
                    h = re.search(r'<a:ext cx="\d+" cy="(\d+)"', sp)
                    if h and int(h.group(1)) >= 365760:
                        rr += 1
            if rr:
                fail(f"p{idx}: roundRect の大きなボックス ×{rr}（直角 rect に）")
            # 塗りありボックスに枠線（大きな図形のみ）
            fb = 0
            for m in re.finditer(r"<p:sp>.*?</p:sp>", xml, re.S):
                sp = m.group(0)
                spPr = re.search(r"<p:spPr>.*?</p:spPr>", sp, re.S)
                if not spPr:
                    continue
                pr = spPr.group(0)
                h = re.search(r'<a:ext cx="\d+" cy="(\d+)"', pr)
                if not (h and int(h.group(1)) >= 365760):
                    continue
                ln = re.search(r"<a:ln[ >].*?</a:ln>", pr, re.S)
                filled = "<a:solidFill>" in re.sub(r"<a:ln[ >].*?</a:ln>", "", pr, flags=re.S)
                if filled and ln and "<a:solidFill>" in ln.group(0):
                    fb += 1
            if fb:
                warn(f"p{idx}: 塗りあり図形に枠線 ×{fb}（塗りカードは line なし — slide-rules §5.3）")
            for c in OLD_COLORS:
                if f'val="{c}"' in xml or f'val="{c.lower()}"' in xml:
                    fail(f"p{idx}: 禁止色 {c}")
        # theme / master も色チェック
        for n in z.namelist():
            if "theme" in n or "slideMaster" in n or "slideLayout" in n:
                xml = z.read(n).decode("utf8", "ignore")
                for c in OLD_COLORS:
                    if f'val="{c}"' in xml:
                        warn(f"{n}: 禁止色 {c}")

    for i, s in enumerate(prs.slides, 1):
        title = ""
        if s.shapes.title is not None and s.shapes.title.has_text_frame:
            title = s.shapes.title.text_frame.text
        else:
            # タイトルPH が無いビルダー: 上部 y<1.2in の最大フォントテキストをタイトルとみなす
            cands = []
            for sh in s.shapes:
                if sh.has_text_frame and sh.top is not None and sh.top < Emu(1097280):
                    sz = max((r.font.size.pt for p in sh.text_frame.paragraphs for r in p.runs if r.font.size), default=0)
                    cands.append((sz, sh.text_frame.text))
            if cands:
                title = max(cands)[1]
        title = title.replace("\n", " ")
        titles.append(title)
        check_title(i, title)
        # サブタイトル疑い: タイトル直下 (y 1.2〜1.75in) の細字テキスト1行
        for sh in s.shapes:
            if sh.has_text_frame and sh.top is not None and Emu(1097280) <= sh.top < Emu(1600200):
                txt = sh.text_frame.text.strip()
                if txt and "\n" not in txt and len(txt) < 60 and txt != title:
                    szs = [r.font.size.pt for p in sh.text_frame.paragraphs for r in p.runs if r.font.size]
                    if szs and max(szs) <= 12 and not any(r.font.bold for p in sh.text_frame.paragraphs for r in p.runs):
                        warn(f"p{i}: タイトル直下にサブタイトルらしき行: 「{txt}」")
    check_terms(term_pages)
    check_ai_smell(term_pages)
    return titles


# ---------------------------------------------------------------- HTML
def check_html(path):
    global STRICT_LEN
    STRICT_LEN = False
    html = Path(path).read_text(encoding="utf8", errors="ignore")
    if not ("338.67mm" in html and "190.5mm" in html):
        fail("16:9 サイズ（338.67mm×190.5mm）が CSS に見当たらない（A4 / 297×167 は旧仕様）")
    if re.search(r"@page\s*{[^}]*297mm", html):
        fail("@page が A4 横のまま")
    for c in OLD_COLORS:
        if re.search(c, html, re.I):
            fail(f"禁止色 #{c}")
    # 角丸
    for m in re.finditer(r"([^{}]{0,80}){[^}]*?border-radius\s*:\s*(\d+(?:\.\d+)?)(px|mm|rem|em)", html):
        sel, v, u = m.group(1), float(m.group(2)), m.group(3)
        if re.search(r"pill|chip|tag|badge|dot", sel, re.I):
            continue  # 小ピルは丸可
        px = v * {"px": 1, "mm": 3.78, "rem": 16, "em": 16}[u]
        if px >= 4:
            fail(f"border-radius {v}{u} on `{sel.strip()[-40:]}`（角丸禁止。小ピル以外は直角）")
            break
    # サブタイトル・図表ラベル
    for cls in ["figttl", "subtitle"]:
        if re.search(r'class="[^"]*\b' + cls + r'\b', html):
            fail(f"サブタイトル/図表ラベル系クラス `.{cls}` が残っている")
    for cls in ["sub", "lead", "caption"]:
        n = len(re.findall(r'class="[^"]*\b' + cls + r'\b', html))
        if n:
            warn(f"`.{cls}` ×{n} — 表紙サブタイトルなら可。コンテンツスライドのタイトル直下なら禁止（目視確認）")
    if re.search(r'<span class="ac">', html):
        fail("タイトル内の色分け <span class=\"ac\"> が残っている")
    # 表ヘッダー
    th = re.search(r"\bth\s*{[^}]*font-size\s*:\s*(\d+)px", html)
    td = re.search(r"\btd\s*{[^}]*font-size\s*:\s*(\d+)px", html)
    if th and td and int(th.group(1)) < int(td.group(1)) + 2:
        fail(f"表ヘッダー {th.group(1)}px が本文 {td.group(1)}px +2pt 未満")
    if re.search(r"\bth\s*{[^}]*font-weight\s*:\s*(400|normal|300)", html):
        fail("表ヘッダーが細字")
    if re.search(r"tr:nth-child\((even|odd)\)", html):
        fail("ゼブラ縞が残っている")
    # 枠線ルール（slide-rules §5.3-5.4 / §6）
    if re.search(r"\btd\b[^{}]*{[^}]*border-bottom\s*:", html) and not re.search(
            r"tr:last-child[^{}]*{[^}]*border(-bottom)?\s*:\s*(0|none)", html):
        fail("最終行の罫線が消えていない（`tr:last-child td{border-bottom:0}` を追加 — 行き先のない罫線禁止）")
    for m in re.finditer(r"([^{}]{0,80}){([^}]*)}", html):
        sel, body = m.group(1), m.group(2)
        if re.search(r"pill|chip|tag|badge|dot|legend", sel, re.I):
            continue
        has_fill = re.search(r"background(-color)?\s*:\s*(?!none|transparent)#?\w", body)
        has_border = re.search(r"border\s*:\s*(?!0|none)\d", body)
        if has_fill and has_border and re.search(r"card|box|pillar|mem|step|stat", sel, re.I):
            warn(f"塗りありボックスに枠線: `{sel.strip()[-40:]}`（塗りカードは border:0 — slide-rules §5.3）")
    # タイトル抽出
    titles = []
    pat = r'<(?:section|div)[^>]*class="(?:[^"]*\bslide\b[^"]*|s|s [^"]*)"[^>]*>'
    parts = re.split(pat, html)
    slides = parts[1:] if len(parts) > 1 else []
    for i, s in enumerate(slides, 1):
        m = re.search(r"<h1[^>]*>(.*?)</h1>", s, re.S) or re.search(r'class="[^"]*\b(?:ttl|title|msg)\b[^"]*"[^>]*>(.*?)</', s, re.S)
        t = re.sub(r"<[^>]+>", "", m.group(1)).strip() if m else ""
        t = re.sub(r"\s+", " ", t)
        titles.append(t)
        check_title(i, t)
    if not slides:
        warn("`.slide` 要素が見つからない（タイトル検査スキップ）")
    if slides:
        check_terms([(i, re.sub(r"<[^>]+>", " ", s)) for i, s in enumerate(slides, 1)])
        check_ai_smell([(i, re.sub(r"<[^>]+>", " ", s)) for i, s in enumerate(slides, 1)])
    else:
        check_terms([(1, re.sub(r"<[^>]+>", " ", html))])
        check_ai_smell([(1, re.sub(r"<[^>]+>", " ", html))])
    return titles


def main():
    if len(sys.argv) < 2:
        sys.exit(__doc__)
    p = sys.argv[1]
    titles = check_pptx(p) if p.lower().endswith(".pptx") else check_html(p)
    print("=== タイトル一覧（上から通し読みしてストーリーが繋がるか確認） ===")
    for i, t in enumerate(titles, 1):
        print(f"{i:>3}  {t or '(なし)'}")
    print()
    for w in warns:
        print("WARN ", w)
    for f in fails:
        print("FAIL ", f)
    print(f"\n{len(fails)} FAIL / {len(warns)} WARN")
    sys.exit(1 if fails else 0)


if __name__ == "__main__":
    main()
