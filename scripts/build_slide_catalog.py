#!/usr/bin/env python3
"""統合カタログ assets/SlideCatalog_16x9.pdf（62型）を組み立てる。

  python3 scripts/build_slide_catalog.py [出力先ディレクトリ=assets]

前提: pipeline/ で `npm run setup` 済み、`pip3 install pymupdf`、LibreOffice（soffice）が PATH にあること。
ghostscript（gs）があれば最後に圧縮する（無ければそのまま出力）。

カタログ（見る側）は1本、生成経路（作る側）は2つ。
  PPTX 経路: pipeline/slide-spec/super_template.json → pipeline/scripts/export_spec_to_editable_pptx.mjs（36型。表紙はHTML側を使うので35型を収録）
  HTML 経路: templates/freeform_parts_16x9.html（27パーツ全部を収録）
両方を 960x540pt に揃えて章立てで並べ替え、通しページ番号と経路ラベルを刷る。
あわせて assets/SuperTemplate_36type.pptx / .pdf と assets/FreeformParts_16x9.pdf も更新する。

外枠（ヘッダーバー・明朝タイトル・フッター）は両経路で揃えてある。
ずれたら export_spec_to_editable_pptx.mjs の addKicker/addTitle/addFooter と
freeform_parts_16x9.html の .bar/h1/.foot を突き合わせること。
"""
import json, math, os, re, shutil, subprocess, sys, tempfile

ROOT = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
PIPE = os.path.join(ROOT, "pipeline")
DECK_NAME = "スライド型カタログ"
OUT_DIR = sys.argv[1] if len(sys.argv) > 1 else os.path.join(ROOT, "assets")

# 日本語フォントを LibreOffice に見せる（macOS。素のままだと日本語が全部落ちる）。
# Linux で使うときは <dir> を自環境のフォントディレクトリに書き換える。
FONTS_CONF = """<?xml version="1.0"?><!DOCTYPE fontconfig SYSTEM "fonts.dtd"><fontconfig>
<dir>/Library/Fonts</dir><dir>/System/Library/Fonts</dir><dir>/System/Library/Fonts/Supplemental</dir>
<dir>~/Library/Fonts</dir><dir>/System/Library/AssetsV2</dir><cachedir>{cache}</cachedir>
<!-- 游ゴシック/游明朝が入っている前提（Office 由来＋macOS 追加書体）。
     ヒラギノに置換すると HTML 側（Chrome＝本物の游書体）と書体が食い違うので、置換せず別名だけ張り、無いときだけヒラギノに落とす。 -->
<alias binding="strong"><family>Yu Gothic</family><prefer><family>Yu Gothic</family><family>YuGothic</family><family>Hiragino Sans</family></prefer></alias>
<alias binding="strong"><family>Yu Mincho</family><prefer><family>Yu Mincho</family><family>YuMincho</family><family>Hiragino Mincho ProN</family></prefer></alias>
<alias binding="strong"><family>Yu Mincho Demibold</family><prefer><family>Yu Mincho Demibold</family><family>YuMincho</family><family>Yu Mincho</family><family>Hiragino Mincho ProN</family></prefer></alias>
</fontconfig>"""

S, P = "S", "P"  # S=PPTX経路(36型) / P=HTML経路(パーツ集)
CATS = [
    ("A", "枠組み", "資料の器をつくる（表紙・全体像・目次・章扉・土台・裏表紙）"),
    ("B", "数字で証明", "1枚1図で主張を裏づける"),
    ("C", "比較・評価", "選択肢や対象を軸で並べて優劣・強弱を見せる"),
    ("D", "構造で通す", "データがない主張を構造で説明する"),
    ("E", "進め方・計画", "段階・時間軸・因果で「どう進むか」を示す"),
    ("F", "提示・締め", "テーマの列挙、外部根拠、意思決定"),
]
ORDER = {
    "A": [(P,1,'表紙'),(P,2,'全体マップ'),(P,3,'目次'),(P,4,'章扉'),(P,27,'セパレーター（アジェンダ再掲）'),
          (S,'executive_summary','エグゼクティブサマリー'),(S,'evidence_basis','調査の土台'),(P,25,'調査の土台'),(P,10,'裏表紙')],
    "B": [(S,'chart_insight','単一チャート＋含意'),(S,'stacked_bar','積み上げ棒'),(S,'waterfall','寄与度ブリッジ'),(S,'true_waterfall','増減ブリッジ（厳密）'),
          (S,'small_multiples','小図の並列比較'),(P,24,'シナリオ線＋成長率チップ'),(P,23,'増減の縦棒＋左右合計'),
          (P,22,'比例円の対比'),(P,17,'分布の順位棒＋注記'),(P,18,'注記つき散布図'),
          (S,'big_stat_pair','大型数値の対比'),(S,'kpi_dashboard','KPI一覧'),(P,7,'大型数値の表＋読み取り')],
    "C": [(S,'comparison_table','選択肢の比較表'),(S,'scenario_table','シナリオ比較表'),(S,'risk_table','リスク一覧表'),(S,'horizontal_axis_table','横軸評価表'),
          (S,'heatmap_table','ヒートマップ表'),(S,'matrix_2x2','2×2マトリクス'),(P,9,'軸のある表'),
          (P,14,'充足度評価表（ハーベイボール）'),(P,13,'状態ヒートマップ＋右コメント'),
          (P,12,'打ち手の効果表'),(P,16,'進捗バブル行列'),(P,15,'割合のドットマトリクス')],
    "D": [(P,11,'主張パネル＋図'),(S,'process_matrix','プロセス×観点の行列'),(S,'nested_row_matrix','入れ子行の行列'),(S,'timeline_matrix','時系列マトリクス'),
          (S,'issue_tree','イシューツリー'),(S,'scr','Situation・Complication・Resolution'),(S,'issue_to_solution_map','課題と打ち手の対応'),
          (P,26,'課題と打ち手の2カラム'),(S,'issue_cause_solution','課題→原因→解決'),(S,'current_target_state','現状と目指す姿'),
          (S,'calc_flow','計算ロジックの流れ'),(P,19,'柱＋土台'),(P,20,'対向シェブロン')],
    "E": [(S,'process_flow','プロセスの段階'),(S,'cycle','循環サイクル'),(S,'chevron_rail','矢羽の段階'),(P,5,'矢羽（プロセス・変遷）'),
          (S,'chevron_value_chain','バリューチェーン'),(S,'decision_fork','分岐と判断'),(P,6,'前提→帰結の2カラム'),
          (S,'roadmap','ロードマップ'),(S,'gantt','ガントチャート')],
    "F": [(S,'theme_card_grid','テーマカード'),(S,'recommendation_pillars','提言の柱'),(S,'numbered_imperatives','番号つき打ち手'),
          (P,8,'並列カード 2×2'),(P,21,'外部動向の根拠グリッド'),(S,'decision_page','意思決定ページ')],
}
LANE = {S: "PPTX", P: "HTML"}
# カタログは HTML パーツ集（ブラウン系）と配色を揃える。SlideSpec に palette が無いときだけ注入する。
# 生成デッキの既定色（ネイビー）は変えない。自社色にするときは SlideSpec ルートの palette を使う。
CATALOG_PALETTE = {"ink": "322014", "navy": "322014", "blue": "5A3921", "cyan": "C5A681", "muted": "8A7B6B",
                   "hair": "E2DCD2", "rose": "A22727", "softBlue": "EFEEE8", "green": "5A3921", "warning": "C5A681"}


def run(cmd, **kw):
    subprocess.run(cmd, check=True, stdout=subprocess.DEVNULL, stderr=subprocess.DEVNULL, **kw)


def resolve_pptx_pages():
    """型ID → PPTX の何ページ目か。スライドを増減してもカタログがずれないようにする。"""
    spec = json.load(open(os.path.join(PIPE, "slide-spec", "super_template.json"), encoding="utf-8"))
    return {sl.get("template"): i for i, sl in enumerate(spec["slides"], start=1)}


def build_plan(pmap):
    plan, page = [], 2  # 1=表紙, 2=索引
    for code, name, desc in CATS:
        page += 1
        plan.append(("divider", code, name, desc, page))
        for src, key, tname in ORDER[code]:
            pno = pmap[key] if src == S else key
            page += 1
            plan.append(("page", src, pno, tname, LANE[src], page, code))
    return plan


def shell_html(plan, style):
    rows = [(e[3], e[4], e[5], e[6]) for e in plan if e[0] == "page"]  # name, lane, pg, cat
    per = math.ceil(len(rows) / 3)
    cols = [rows[i * per:(i + 1) * per] for i in range(3)]

    def col(c):
        return "\n".join(
            f'<div class="ix"><span class="ixc">{cat}</span><span class="ixn">{n}</span>'
            f'<span class="ixl">{l}</span><span class="ixp">P.{pg}</span></div>'
            for n, l, pg, cat in c)

    divs = "\n".join(f'''<section class="s chap">
  <div class="bar"><div class="logo">{DECK_NAME}</div><div class="date">{e[1]}</div></div>
  <div class="cno">{e[1]}</div><div class="cbig">{e[2]}</div><div class="cdesc">{e[3]}</div>
  <div class="foot"><b>{DECK_NAME}</b><span>{e[4]}</span></div>
</section>''' for e in plan if e[0] == "divider")

    extra = '''<style>
.cdesc{font-size:12pt;color:#6b4a2b;margin-top:4mm;line-height:1.7}
.ixwrap{display:grid;grid-template-columns:repeat(3,1fr);gap:0 7mm;flex:1;min-height:0;align-content:start}
.ix{display:grid;grid-template-columns:5mm 1fr 11mm 11mm;gap:1.6mm;align-items:baseline;padding:.85mm 0;border-bottom:1px solid #ece6db;font-size:8.2pt}
.ixc{font-family:"Yu Mincho",serif;font-weight:700;color:#6b4a2b}
.ixn{color:#241a10}.ixl{font-size:7.2pt;color:#8a7c6c;text-align:right}
.ixp{font-size:7.6pt;color:#8a7c6c;text-align:right}
.cover .lead{font-size:11pt;color:#4a3d30;line-height:1.9;margin-top:6mm;max-width:180mm}
</style>'''
    return f'''<!DOCTYPE html><html lang="ja"><head><meta charset="utf-8">
<title>{DECK_NAME} 16:9</title>{style}{extra}</head><body>
<section class="s cover">
  <div class="bar"><div class="logo">{DECK_NAME}</div><div class="date">16:9 ｜ {len(rows)}型</div></div>
  <div class="big">{DECK_NAME}</div><div class="rule"></div>
  <div class="lead">各ページのタイトル欄は<b>型名</b>。見本の主張文は置いていない — 実デッキではストーリーラインから起こした主張文に必ず差し替える（slide-rules §2.8）。<br>
  右上の「PPTX」「HTML」はその型を出せる経路。PPTX＝36型パイプライン（編集可能PPTXで書き出す）、HTML＝自由記述パーツ集。</div>
  <div class="foot"><b>{DECK_NAME}</b><span>1</span></div>
</section>
<section class="s">
  <div class="bar"><div class="logo">{DECK_NAME}</div><div class="date">索引</div></div>
  <h1>収録している型</h1>
  <div class="c"><div class="ixwrap">
      <div>{col(cols[0])}</div>
      <div>{col(cols[1])}</div>
      <div>{col(cols[2])}</div>
  </div></div>
  <div class="foot"><b>{DECK_NAME}</b><span>2</span></div>
</section>
{divs}
</body></html>'''


def main():
    import fitz
    tmp = tempfile.mkdtemp(prefix="slidecatalog-")
    cache = os.path.join(tmp, "fccache"); os.makedirs(cache, exist_ok=True)
    conf = os.path.join(tmp, "fonts.conf")
    open(conf, "w").write(FONTS_CONF.format(cache=cache))
    plan = build_plan(resolve_pptx_pages())

    parts_src = os.path.join(ROOT, "templates", "freeform_parts_16x9.html")
    html = open(parts_src, encoding="utf-8").read()
    style = html[html.index("<style>"):html.index("</style>") + len("</style>")]

    # パーツ集はテンプレなので「◯◯」が入っている。カタログではPPTX側と同じ資料名に差し替える
    cat_parts = os.path.join(tmp, "parts_catalog.html")
    open(cat_parts, "w", encoding="utf-8").write(
        html.replace('<div class="logo">資料名</div>', f'<div class="logo">{DECK_NAME}</div>')
            .replace("<b>資料名</b>", f"<b>{DECK_NAME}</b>"))

    open(os.path.join(tmp, "shell.html"), "w", encoding="utf-8").write(shell_html(plan, style))

    h2p = os.path.join(PIPE, "scripts", "html_to_pdf.mjs")
    run(["node", h2p, cat_parts, os.path.join(tmp, "parts.pdf")], cwd=PIPE)
    # 単体配布用は「◯◯」プレースホルダーのまま出す（テンプレとして渡すもの）
    run(["node", h2p, parts_src, os.path.join(tmp, "parts_plain.pdf")], cwd=PIPE)
    run(["node", h2p, os.path.join(tmp, "shell.html"), os.path.join(tmp, "shell.pdf")], cwd=PIPE)

    spec = os.path.join(PIPE, "slide-spec", "super_template.json")
    spec_obj = json.load(open(spec, encoding="utf-8"))
    if "palette" not in spec_obj:
        spec_obj["palette"] = CATALOG_PALETTE
        spec = os.path.join(tmp, "super_template.catalog.json")
        json.dump(spec_obj, open(spec, "w", encoding="utf-8"), ensure_ascii=False, indent=2)
    pptx = os.path.join(tmp, "SuperTemplate_36type.pptx")
    run(["node", os.path.join(PIPE, "scripts", "export_spec_to_editable_pptx.mjs"), spec, pptx], cwd=PIPE)
    env = dict(os.environ, FONTCONFIG_FILE=conf)
    run(["soffice", "--headless", "--convert-to", "pdf", "--outdir", tmp, pptx], env=env)
    spdf = os.path.join(tmp, os.path.basename(pptx).replace(".pptx", ".pdf"))

    shell, Sdoc, Pdoc = (fitz.open(os.path.join(tmp, "shell.pdf")), fitz.open(spdf),
                         fitz.open(os.path.join(tmp, "parts.pdf")))
    out = fitz.open()
    out.insert_pdf(shell, from_page=0, to_page=1)
    di, lanes = 2, {}
    for e in plan:
        if e[0] == "divider":
            out.insert_pdf(shell, from_page=di, to_page=di); di += 1
        else:
            _, src, pno, _n, lane, pg, _c = e
            out.insert_pdf(Sdoc if src == S else Pdoc, from_page=pno - 1, to_page=pno - 1)
            lanes[pg] = lane
    for i, page in enumerate(out, start=1):
        if i <= 2:
            continue
        # 元のページ番号を伏せて通し番号に刷り直す。フッター文字のベースラインは
        # PPTX側 FOOTER_Y+0.09in・HTML側 .foot ともに約 524pt に揃えてある。
        page.draw_rect(fitz.Rect(828, 505, 952, 537), color=None, fill=(1, 1, 1))
        page.insert_text((938 - len(str(i)) * 5, 524), str(i), fontsize=8.5,
                         fontname="helv", color=(0.54, 0.49, 0.42))
        if i in lanes:
            page.insert_text((876, 524), lanes[i], fontsize=7.5, fontname="helv",
                             color=(0.62, 0.57, 0.50))
    out.set_metadata({"title": f"{DECK_NAME} 16:9"})
    raw = os.path.join(tmp, "raw.pdf"); out.save(raw)

    os.makedirs(OUT_DIR, exist_ok=True)
    final = os.path.join(OUT_DIR, "SlideCatalog_16x9.pdf")
    # ページごとにフォントが重複埋め込みされて 16MB 級になるので圧縮する
    try:
        run(["gs", "-sDEVICE=pdfwrite", "-dCompatibilityLevel=1.5", "-dPDFSETTINGS=/prepress",
             "-dNOPAUSE", "-dQUIET", "-dBATCH", "-dSubsetFonts=true",
             f"-sOutputFile={final}", raw])
    except Exception:
        shutil.copy(raw, final)
    # 単体の一覧もあわせて更新
    shutil.copy(os.path.join(tmp, "parts_plain.pdf"), os.path.join(OUT_DIR, "FreeformParts_16x9.pdf"))
    shutil.copy(spdf, os.path.join(OUT_DIR, "SuperTemplate_36type.pdf"))
    dst_pptx = os.path.join(OUT_DIR, "SuperTemplate_36type.pptx")
    if os.path.abspath(dst_pptx) != os.path.abspath(pptx):
        shutil.copy(pptx, dst_pptx)
    n62 = build_62_pptx(raw, pptx, plan, os.path.join(OUT_DIR, "SuperTemplate_62type.pptx"))
    print(f"{final}  ({out.page_count}p, {os.path.getsize(final)/1e6:.1f}MB)")
    print(f"{os.path.join(OUT_DIR, 'SuperTemplate_62type.pptx')}  ({n62} slides: 35 native + {n62-35} image)")
    shutil.rmtree(tmp, ignore_errors=True)


def build_62_pptx(raw_pdf, pptx36, plan, out_path, dpi=192):
    """62型を1本にした PPTX。SlideSpec由来の35型はネイティブ編集可能図形のまま、HTML経路の27パーツと
    索引・章扉はカタログPDFのページを画像で収録する（HTML側はSlideSpec外なので図形化できない）。
    スライド番号はカタログPDFのページ番号と一致させ、フッターの番号も通し番号に書き換える。"""
    import io
    import fitz
    from pptx import Presentation
    from pptx.util import Pt, Inches
    from pptx.enum.text import PP_ALIGN

    doc = fitz.open(raw_pdf)
    prs = Presentation(pptx36)
    W, H = prs.slide_width, prs.slide_height
    slides36 = list(prs.slides)
    sld_ids = list(prs.slides._sldIdLst)
    layout = next((l for l in prs.slide_layouts if len(l.placeholders) == 0), prs.slide_layouts[0])

    pages = {1: ("img", 0, "表紙"), 2: ("img", 1, "索引")}
    for e in plan:
        if e[0] == "divider":
            pages[e[4]] = ("img", e[4] - 1, f"章扉 {e[1]} {e[2]}")
        else:
            _, src, pno, tname, lane, pg, _c = e
            pages[pg] = ("S", pno, tname) if src == S else ("img", pg - 1, f"パーツ{pno:02d} {tname}")

    order = []
    for pg in sorted(pages):
        kind, v, label = pages[pg]
        if kind == "S":
            slide = slides36[v - 1]
            # フッターの番号（右下・元は 1..36）を通し番号に。隣に経路ラベル PPTX を刷る
            for sh in slide.shapes:
                if sh.has_text_frame and sh.text_frame.text.strip() == str(v) and sh.left > W * 0.8:
                    run = sh.text_frame.paragraphs[0].runs[0]
                    run.text = str(pg)
                    tb = slide.shapes.add_textbox(sh.left - Inches(1.0), sh.top, Inches(0.95), sh.height)
                    tf = tb.text_frame; tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
                    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.RIGHT
                    r = p.add_run(); r.text = "PPTX"; r.font.size = Pt(7.5); r.font.name = run.font.name
                    try:
                        r.font.color.rgb = run.font.color.rgb
                    except Exception:
                        pass
                    break
            order.append(sld_ids[v - 1])
        else:
            s = prs.slides.add_slide(layout)
            for ph in list(s.placeholders):
                ph._element.getparent().remove(ph._element)
            png = doc[v].get_pixmap(dpi=dpi).tobytes("png")
            s.shapes.add_picture(io.BytesIO(png), 0, 0, W, H)
            s.notes_slide.notes_text_frame.text = (
                f"{label}：HTML経路（templates/freeform_parts_16x9.html）のパーツを画像で収録。"
                "編集する場合はHTML側のパーツをコピーして組む。" if label.startswith("パーツ")
                else f"{label}：カタログの構成ページ（画像）。")
            order.append(prs.slides._sldIdLst[-1])

    lst = prs.slides._sldIdLst
    for el in list(lst):
        lst.remove(el)
    for el in order:
        lst.append(el)
    prs.save(out_path)
    return len(order)


if __name__ == "__main__":
    main()
